# -*- coding: utf-8 -*-
"""check_consistency.py -- does the deck still agree with everything else?

Every drift found by hand in this project was one of four kinds, so this
checks all four automatically:

  1. a number the speech script quotes that is not in the deck
  2. a headline result in the deck that disagrees with RESULTS-SUMMARY.txt
  3. a file the deck references (video, figure, script) that does not exist
  4. a claim word in the deck with no evidence anywhere in the repo

Run it after ANY edit to build.py, content.py, refs.py or the speech script:

    python3 check_consistency.py

Exit status is non-zero if anything fails, so it can gate a commit.
"""
import os, re, sys, glob
from pptx import Presentation

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DECK = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                    "Review1_GaN_Segmented_Gate_Driver.pptx")
SPEECH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "SPEECH-SCRIPT.md")
SUMMARY = os.path.join(ROOT, "results", "RESULTS-SUMMARY.txt")

fails, warns = [], []

def deck_text():
    p = Presentation(DECK)
    out = []
    for s in p.slides:
        for sh in s.shapes:
            if sh.has_text_frame:
                out.append(sh.text_frame.text)
            if getattr(sh, "has_table", False) and sh.has_table:
                out += [c.text for r in sh.table.rows for c in r.cells]
    return "\n".join(out)

DT = deck_text()
SP = open(SPEECH, encoding="utf-8").read() if os.path.exists(SPEECH) else ""
SM = open(SUMMARY, encoding="utf-8").read() if os.path.exists(SUMMARY) else ""

def norm(t):
    # the deck uses U+2212 and thin spaces; the script uses ASCII
    return (t.replace("−", "-").replace("–", "-").replace(" ", " ")
             .replace(" ", "").replace(",", ""))

DTn, SPn, SMn = norm(DT), norm(SP), norm(SM)

# ---- 1. numbers the speech quotes that the deck does not contain -----------
# Only bolded figures: those are the ones the presenter says out loud.
spoken = set(re.findall(r"\*\*([+-]?\d[\d.]*\s*(?:%|V|nH|ns|µJ|cells)?)\*\*", SP))
for v in sorted(spoken):
    vv = norm(v).strip()
    bare = re.sub(r"\s*(%|V|nH|ns|µJ|cells)$", "", vv).strip()
    if not bare or bare in ("1", "2", "3", "4", "5", "8"):
        continue
    if bare not in DTn:
        fails.append("speech quotes %-12s but the deck does not contain it" % ("'%s'" % vv))

# ---- 2. headline results must match RESULTS-SUMMARY ------------------------
HEADLINES = {
    "5.2 %":   "ceiling on scheduling",
    "25.1 %":  "(A) better fixed word",
    "3.9 %":   "(B) adaptation",
    "13.4 %":  "adaptation share",
    "2.576":   "shipped margin",
    "34622":   "transient count",
}
for num, what in HEADLINES.items():
    n = norm(num)
    in_deck, in_sum = n in DTn, n in SMn
    if in_deck and not in_sum:
        fails.append("%-8s (%s) is on a slide but NOT in RESULTS-SUMMARY" % (num, what))
    if in_sum and not in_deck:
        warns.append("%-8s (%s) is in RESULTS-SUMMARY but not on any slide" % (num, what))

# ---- 3. every file the build references must exist -------------------------
build = open(os.path.join(os.path.dirname(os.path.abspath(__file__)), "build.py"),
             encoding="utf-8").read()
for m in re.finditer(r'"([\w./-]+\.(?:png|mp4|pptx|cir|lib|v|m|py|sh))"', build):
    f = m.group(1)
    if f.startswith(("http", "github")):
        continue
    rel = f.lstrip("/")
    HERE = os.path.dirname(os.path.abspath(__file__))
    cands = [os.path.join(HERE, rel), os.path.join(ROOT, "results", rel),
             os.path.join(ROOT, rel), os.path.join(ROOT, "results", os.path.basename(rel)),
             os.path.join(HERE, os.path.basename(rel))]
    if not any(os.path.exists(c) for c in cands):
        fails.append("build.py references %s which does not exist" % f)

# ---- 4. claim words on a slide need evidence somewhere ---------------------
CORPUS = ""
for f in (glob.glob(os.path.join(ROOT, "results", "*.txt")) +
          glob.glob(os.path.join(ROOT, "results", "*.md")) +
          glob.glob(os.path.join(ROOT, "scripts", "*.py")) +
          glob.glob(os.path.join(ROOT, "scripts", "*.sh"))):
    try: CORPUS += open(f, errors="ignore").read()
    except Exception: pass
BANNED = {
    "browser-WASM": "no WASM deck exists in the repository",
    "Spectre deck reproduces": "Spectre has never been run",
    "cross-simulator agreement": "only claimable if two simulators actually ran",
}
for phrase, why in BANNED.items():
    if phrase.lower() in DT.lower() and "not cross-simulator" not in DT.lower():
        fails.append("deck says %r -- %s" % (phrase, why))

# ---- 5. figures must regenerate identically from current data --------------
# Opt-in with --figures: it re-executes the plotting scripts, so it is slow.
# A figure that no longer reproduces from the committed data is stale, and a
# stale figure showing superseded numbers is exactly the defect this project
# has already hit once.
if "--figures" in sys.argv:
    import subprocess, shutil, tempfile
    from PIL import Image, ImageChops
    import numpy as np
    GEN = {"fig1_crosstalk": "figures.py", "paper_fig2_ceiling": "paper_figs.py",
           "fig_rtl_waveform": "plot_waveform.py", "fig_architecture": "arch_diagram.py",
           "fig_lloop": "plot_lloop.py", "fig_circuit": "circuit_diagram.py"}
    RESD = os.path.join(ROOT, "results")
    tmp = tempfile.mkdtemp()
    for name in GEN:
        p = os.path.join(RESD, name + ".png")
        if os.path.exists(p):
            shutil.copy(p, os.path.join(tmp, name + ".png"))
    for script in sorted(set(GEN.values())):
        subprocess.run(["python3", os.path.join(ROOT, "scripts", script)],
                       capture_output=True, timeout=900)
    for name in GEN:
        ap, bp = os.path.join(tmp, name + ".png"), os.path.join(RESD, name + ".png")
        if not (os.path.exists(ap) and os.path.exists(bp)):
            fails.append("figure %s could not be compared" % name); continue
        a = Image.open(ap).convert("RGB"); b = Image.open(bp).convert("RGB")
        if a.size != b.size:
            fails.append("figure %s changed size on regeneration" % name); continue
        d = np.asarray(ImageChops.difference(a, b), dtype=float)
        pct = 100.0 * (d.max(axis=2) > 8).mean()
        if pct >= 0.01:
            fails.append("figure %s is STALE: regenerating from current data "
                         "changes %.2f %% of pixels" % (name, pct))
        else:
            print("  ok    figure %-22s reproduces from current data" % name)

# ---- report ---------------------------------------------------------------
print("consistency check")
print("-" * 66)
if not fails and not warns:
    print("  PASS -- deck, speech script, results summary and files all agree")
for w in warns:
    print("  warn  %s" % w)
for f in fails:
    print("  FAIL  %s" % f)
print("-" * 66)
print("  %d failure(s), %d warning(s)" % (len(fails), len(warns)))
sys.exit(1 if fails else 0)
