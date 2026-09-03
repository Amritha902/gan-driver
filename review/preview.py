# -*- coding: utf-8 -*-
"""preview.py -- approximate visual render of the deck, for eyeball QA.

LibreOffice cannot be installed in this container, so slides have never been
seen rendered. This draws each slide from python-pptx geometry with PIL:
text boxes with their real text wrapped at the real box width and font size,
pictures as labelled placeholders, tables as grids.

It is NOT a faithful renderer -- fonts, kerning and autofit differ. What it
DOES show reliably is layout: overlaps, text running past its box, elements
off the slide, uneven gaps. That is what has been unverifiable until now.

    python3 preview.py            # all slides -> preview/slide-NN.png
    python3 preview.py 7 13       # only those
"""
import os, sys
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu

EMU = 914400.0
SCALE = 110                      # px per inch
OUT = "preview"

def font(px, bold=False):
    for p in ("/usr/share/fonts/truetype/dejavu/DejaVuSans%s.ttf"
              % ("-Bold" if bold else ""),
              "/usr/share/fonts/truetype/liberation/LiberationSans%s.ttf"
              % ("-Bold" if bold else "-Regular")):
        if os.path.exists(p):
            try: return ImageFont.truetype(p, max(7, int(px)))
            except Exception: pass
    return ImageFont.load_default()

def wrap(draw, text, f, width):
    out = []
    for para in text.split("\n"):
        if not para.strip():
            out.append("")
            continue
        line = ""
        for w in para.split(" "):
            t = (line + " " + w).strip()
            if draw.textlength(t, font=f) <= width or not line:
                line = t
            else:
                out.append(line); line = w
        out.append(line)
    return out

def runs_of(shape):
    """Largest font size on the shape, and whether its first run is bold."""
    sz, bold = None, False
    try:
        for para in shape.text_frame.paragraphs:
            for r in para.runs:
                if r.font.size:
                    s = r.font.size.pt
                    sz = s if sz is None else max(sz, s)
                if r.font.bold and bold is False:
                    bold = True
    except Exception:
        pass
    return sz or 14.0, bold

OVER = []

def render(slide, n, W, H):
    img = Image.new("RGB", (int(W * SCALE), int(H * SCALE)), "white")
    d = ImageDraw.Draw(img)
    for sh in slide.shapes:
        if sh.left is None or sh.top is None or not sh.width or not sh.height:
            continue
        x, y = sh.left / EMU * SCALE, sh.top / EMU * SCALE
        w, h = sh.width / EMU * SCALE, sh.height / EMU * SCALE
        st = str(sh.shape_type)
        if "PICTURE" in st:
            d.rectangle([x, y, x + w, y + h], outline=(150, 150, 150), width=2)
            d.line([x, y, x + w, y + h], fill=(205, 205, 205))
            d.line([x, y + h, x + w, y], fill=(205, 205, 205))
            d.text((x + 6, y + 6), "IMAGE", font=font(11), fill=(120, 120, 120))
            continue
        if getattr(sh, "has_table", False) and sh.has_table:
            t = sh.table
            rows, cols = len(t.rows), len(t.columns)
            ch = h / max(rows, 1)
            cw = [c.width / EMU * SCALE for c in t.columns]
            cx = x
            for ci in range(cols):
                cy = y
                for ri in range(rows):
                    d.rectangle([cx, cy, cx + cw[ci], cy + ch],
                                outline=(170, 170, 170))
                    txt = t.cell(ri, ci).text
                    if txt:
                        f = font(8, ri == 0)
                        for k, ln in enumerate(wrap(d, txt, f, cw[ci] - 8)[:6]):
                            d.text((cx + 4, cy + 3 + k * 10), ln, font=f,
                                   fill=(20, 20, 20))
                    cy += ch
                cx += cw[ci]
            continue
        if not (sh.has_text_frame and sh.text_frame.text.strip()):
            d.rectangle([x, y, x + w, y + h], outline=(225, 225, 225))
            continue
        # Per-paragraph sizing. Using one size for the whole shape made small
        # captions under a big number render at the number's size, which
        # invented overflow that is not in the file.
        blocks = []
        for para in sh.text_frame.paragraphs:
            txt = "".join(r.text for r in para.runs)
            if not txt.strip():
                blocks.append(("", 11.0, False)); continue
            sz = next((r.font.size.pt for r in para.runs if r.font.size), None) or 14.0
            bd = any(r.font.bold for r in para.runs)
            blocks.append((txt, sz, bd))
        if not blocks:
            blocks = [(sh.text_frame.text, 14.0, False)]
        need = 0.0
        laid = []
        for txt, sz, bd in blocks:
            f = font(sz * SCALE / 72.0 * 0.92, bd)
            lh = (sz * SCALE / 72.0) * 1.22
            lns = wrap(d, txt, f, w - 6) if txt else [""]
            laid.append((lns, f, lh))
            need += len(lns) * lh
        over = need > h + 2
        if over:
            OVER.append((n, (sh.name or "")[:16], round(need/SCALE,2), round(h/SCALE,2),
                         sh.text_frame.text.strip().split("\n")[0][:46]))
        d.rectangle([x, y, x + w, y + h],
                    outline=(220, 60, 60) if over else (215, 215, 215),
                    width=2 if over else 1)
        cy = y
        for lns, f, lh in laid:
            for ln in lns:
                if cy > y + h + lh * 3:
                    break
                d.text((x + 3, cy), ln, font=f,
                       fill=(200, 30, 30) if cy > y + h else (15, 15, 15))
                cy += lh
    d.rectangle([0, 0, img.width - 1, img.height - 1], outline=(120, 120, 120))
    return img

def main():
    prs = Presentation("Review1_GaN_Segmented_Gate_Driver.pptx")
    W = prs.slide_width / EMU
    H = prs.slide_height / EMU
    os.makedirs(OUT, exist_ok=True)
    want = [int(a) for a in sys.argv[1:]] or None
    for i, s in enumerate(prs.slides, 1):
        if want and i not in want:
            continue
        p = os.path.join(OUT, "slide-%02d.png" % i)
        render(s, i, W, H).save(p)
    if OVER:
        print("\nTEXT RUNNING PAST ITS BOX:")
        for n, nm, need, box, txt in OVER:
            print("  s%-3d %-16s needs %.2f in of %.2f in  | %s" % (n, nm, need, box, txt))
    else:
        print("\nno text overflow detected")
    print("\n%d slides rendered into %s/" % (len(prs.slides) if not want else len(want), OUT))

if __name__ == "__main__":
    main()
