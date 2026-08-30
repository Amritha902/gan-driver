"""Fill the Review-1 template from the project's actual results.

Edits the template in place via lxml on each slide part, so every inherited
property (theme, fonts, the VIT logo, page numbers) survives untouched. Only
the paragraphs inside the named text boxes are replaced.
"""
import copy, re
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt

A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS = {"a": A}
def q(t): return "{%s}%s" % (A, t)

PARA_TPL = (
 '<a:p xmlns:a="%s"><a:pPr marL="{marL}" indent="{indent}">'
 '<a:spcAft><a:spcPts val="{spc}"/></a:spcAft>'
 '<a:buSzPct val="100000"/>{bu}</a:pPr>{runs}</a:p>' % A)
RUN_TPL = (
 '<a:r><a:rPr lang="en-US" sz="{sz}"{b} dirty="0">'
 '<a:solidFill><a:srgbClr val="{col}"/></a:solidFill>'
 '<a:latin typeface="Calibri" pitchFamily="34" charset="0"/>'
 '<a:ea typeface="Calibri" pitchFamily="34" charset="-122"/>'
 '<a:cs typeface="Calibri" pitchFamily="34" charset="-120"/>'
 '</a:rPr><a:t xml:space="preserve">{t}</a:t></a:r>')

def esc(s):
    return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

def para(runs, level=0, sz=1600, spc=800, bullet=True):
    """runs: list of (text, bold) or a plain string."""
    if isinstance(runs, str):
        runs = [(runs, False)]
    marL = [177800, 520700, 863600][level]
    ind = -177800 if bullet else 0
    ch = ["•", "–", "•"][level]
    bu = '<a:buChar char="%s"/>' % ch if bullet else "<a:buNone/>"
    body = "".join(RUN_TPL.format(sz=sz, b=' b="1"' if b else "",
                                  col="000000", t=esc(t)) for t, b in runs)
    return PARA_TPL.format(marL=marL, indent=ind if bullet else 0,
                           spc=spc, bu=bu, runs=body)

def set_body(shape, paras):
    tx = shape.text_frame._txBody
    for p in tx.findall(q("p")):
        tx.remove(p)
    for xml in paras:
        tx.append(etree.fromstring(xml))

def find_shape(slide, needle):
    for sh in slide.shapes:
        if sh.has_text_frame and needle in sh.text_frame.text:
            return sh
    raise KeyError(needle)

def set_title(slide, text):
    sh = None
    for s in slide.shapes:
        if s.has_text_frame and s.width and abs(s.width - Inches(10.5)) < Inches(0.2) \
           and s.top is not None and s.top < Inches(1.0):
            sh = s; break
    if sh is None: return
    for p in sh.text_frame.paragraphs:
        for r in p.runs:
            r.text = text; return
