# -*- coding: utf-8 -*-
import sys, os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from fill import para, set_body, find_shape, set_title, q, A, esc, RUN_TPL, PARA_TPL
import content

RES = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "results")
SRC = "template_ext.pptx"
OUT = "Review1_GaN_Segmented_Gate_Driver.pptx"

# ---- text-fit estimate (no renderer available in this environment) ---------
# Calibri average advance is ~0.478 em for mixed-case prose; 0.50 is used here
# so the estimate errs toward predicting overflow rather than hiding it.
CH_EM, LINE_EM, BOX_W, BOX_H = 0.50, 1.22, 11.90, 5.40
INDENT = {0: 0.194, 1: 0.569}

def fits(paras, spc_pt=6.0):
    h = 0.0
    for runs, lvl in paras:
        n = sum(len(t) for t, _ in runs)
        sz = 16 if lvl == 0 else 15
        cw = sz * CH_EM / 72.0
        avail = BOX_W - INDENT[lvl]
        lines = max(1, -(-int(n * cw / avail * 100) // 100) if n else 1)
        lines = max(1, int(n * cw / avail) + (1 if (n * cw) % avail else 0))
        h += lines * sz * LINE_EM / 72.0 + spc_pt / 72.0
    return h

def build_paras(spec):
    out = []
    for runs, lvl in spec:
        out.append(para(runs, level=lvl, sz=(1600 if lvl == 0 else 1500),
                        spc=600, bullet=True))
    return out

# A previous run's output must not survive a crash: qa.py and validate.py
# read OUT by name, and a stale file from the last SUCCESSFUL build looks
# exactly like a passing result. Delete it before doing any work.
if os.path.exists(OUT):
    os.remove(OUT)

p = Presentation(SRC)
S = p.slides

# ---------------- slide 2 : title -----------------------------------------
def set_run_after(slide, label, newtext):
    """Replace the text of the first run that follows a run equal to `label`."""
    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        runs = [r for para_ in sh.text_frame.paragraphs for r in para_.runs]
        for i, r in enumerate(runs):
            if r.text.strip() == label and i + 1 < len(runs):
                runs[i + 1].text = newtext
                return True
    return False

for sh in S[1].shapes:
    if sh.has_text_frame:
        for pa in sh.text_frame.paragraphs:
            for r in pa.runs:
                if r.text.strip() == "Project Title":
                    r.text = ("A Segmented Gate Driver for GaN HEMTs: Measuring What "
                              "Operating-Point Scheduling Is Actually Worth")
set_run_after(S[1], "Student Name(s):", "Sanjay Kumar 23BEC1447  ·  Aamir Abdullah 23BPS1197  ·  Amritha S 23BEC1368")
set_run_after(S[1], "Guide:", "Dr. Bindu, School of Electronics Engineering (SENSE)")

# ---------------- slides 4, 6, 7 : bulleted content ------------------------
for idx, spec, needle in ((3, content.SLIDE4, "Problem Statement:"),
                          (6, content.SLIDE6, "Proposed Solution:"),
                          (8, content.SLIDE7,  "Work completed so far"),
                          (9, content.SLIDE7B, "Problem Statement:")):
    sh = find_shape(S[idx], needle)
    set_body(sh, build_paras(spec))
    print("slide %d: %.2f in of %.2f  %s"
          % (idx + 1, fits(spec), BOX_H, "OK" if fits(spec) <= BOX_H else "OVERFLOW"))

rtl_shape = find_shape(S[10], "Problem Statement:")
set_body(rtl_shape, build_paras(content.SLIDE_RTL_SHORT))
set_title(S[10], "FPGA Controller — verified RTL")
print("slide 11 (RTL text block): %.2f in of 2.45  %s"
      % (fits(content.SLIDE_RTL_SHORT),
         "OK" if fits(content.SLIDE_RTL_SHORT) <= 2.45 else "OVERFLOW"))

# The official brief says 10 minutes; the template stub said 8-10.
for sh in S[2].shapes:
    if sh.has_text_frame:
        for pa in sh.text_frame.paragraphs:
            for r in pa.runs:
                if "8\u201310 minutes" in r.text:
                    r.text = r.text.replace("8\u201310 minutes", "10 minutes")

set_title(S[8], "Work Completed — 50 %")
set_title(S[9], "Timeline, Milestones & Tools")

p.save(OUT)
print("wrote", OUT)

# ---------------- slides 5 & 6 : literature survey (4 refs each) -----------
from pptx.dml.color import RGBColor
import refs as R

def set_cell(cell, runs, size=Pt(10.5), grey=False):
    """runs: list of (text, bold) or a plain string.

    A newline inside an <a:t> is not a line break in OOXML - PowerPoint
    swallows it - so "\n" starts a new PARAGRAPH instead."""
    if isinstance(runs, str):
        runs = [(runs, False)]
    tf = cell.text_frame
    tf.word_wrap = True
    pa = tf.paragraphs[0]
    for r in list(pa.runs):
        r._r.getparent().remove(r._r)
    for extra in list(tf.paragraphs)[1:]:
        extra._p.getparent().remove(extra._p)
    col = RGBColor(0x7A, 0x7A, 0x7A) if grey else RGBColor(0, 0, 0)
    for txt, bold in runs:
        for i, line in enumerate(txt.split("\n")):
            if i:
                pa = tf.add_paragraph()
            if not line:
                continue
            run = pa.add_run(); run.text = line
            run.font.size = size; run.font.bold = bold; run.font.name = "Calibri"
            run.font.color.rgb = col

def move_note(slide, top_in):
    """A table grows past its declared height when a cell wraps, so the
    footnote has to be placed below where the table ACTUALLY ends, not
    below where python-pptx thinks it does."""
    n = None
    for sh in slide.shapes:
        if sh.has_table or not sh.has_text_frame:
            continue
        # the footnote is the wide low text box, not the title or page number
        if sh.top > Inches(5.0) and sh.width > Inches(6.0):
            n = sh
    if n is not None:
        n.top = Inches(top_in)


def table_of(slide):
    for sh in slide.shapes:
        if sh.has_table: return sh.table
    raise KeyError("no table")

# ---- slide 5: the LANDSCAPE, one slide, grouped by what the work does ----
# The old deck spread eight papers over two identical 4-row tables and made the
# reader do the clustering. Reviewers want the shape of the field on one slide
# and the near neighbours on the next, so that is what these two now are.
CLUSTERS = [
 ("A · Passive / analogue crosstalk fixes",
  "[1] [2] [3] [4] [11] [12]\n[14]–[19] [21]",
  "Negative off-bias, RC-diode rails, three-level drive, Miller clamps. Cheap and effective, "
  "but the setting is chosen once by hand and never revisited."),
 ("B · Closed-loop analogue adaptation",
  "[5] [6] [7]\n[20] [22] [23] [29] [30]",
  "Sense the operating point, regulate drive strength or timing in the loop. Reports large "
  "gains — 30.5 % less overshoot, 75 % less turn-off loss — against a conventional driver."),
 ("C · Adaptive dead-time control",
  "[8] [13]\n[25] [26] [28]",
  "Dead time driven to sub-nanosecond across load. The one field our own freeze test finds "
  "actually carries the benefit — and only at light load."),
 ("D · DIGITAL / segmented drive  — where this project sits",
  "[9] BASE   [10]\n[24] [27]",
  "The gate waveform selected by a multibit CODE rather than a resistor network. [9] is the "
  "BASE PAPER — doi.org/10.1002/cta.3136 — and [10] is the closest GaN implementation. "
  "Digital means FPGA-implementable, and the code space can be searched exhaustively."),
]
sc = S[4]
set_title(sc, "Literature Landscape — four clusters, and where we sit")
note = find_shape(sc, "Minimum 8")
set_body(note, [para([("%d references, all publisher-verified, grouped by what the driver DOES. " % len(R.REFS) + 
                       "Clusters A\u2013C set or regulate in analogue; only cluster D makes the "
                       "setting a digital CODE \u2014 which is what makes an exhaustive search "
                       "possible at all.",
                       False)], level=0, sz=1150, spc=0, bullet=False)])
tbl = table_of(sc)
set_cell(tbl.cell(0, 0), "#");        set_cell(tbl.cell(0, 1), "Cluster")
set_cell(tbl.cell(0, 2), "Refs");     set_cell(tbl.cell(0, 3), "What the cluster does")
set_cell(tbl.cell(0, 4), "")
tbl.cell(0, 3).merge(tbl.cell(0, 4))
for row, (name, refs_, what) in enumerate(CLUSTERS, start=1):
    set_cell(tbl.cell(row, 0), "ABCD"[row - 1])
    set_cell(tbl.cell(row, 1), [(name, True)])
    set_cell(tbl.cell(row, 2), refs_)
    tbl.cell(row, 3).merge(tbl.cell(row, 4))
    set_cell(tbl.cell(row, 3), what)
for row in (6, 5):
    tbl._tbl.remove(tbl.rows[row]._tr)
for row in range(len(tbl.rows)):
    tbl.rows[row].height = Inches(0.90 if row else 0.35)
move_note(sc, 6.45)

# ---- slide 6: the five closest papers, base paper first -----------------
CLOSEST = [R.REFS[8]] + [r for r in R.REFS if r.closest][:4]
sn = S[5]
set_title(sn, "The Five Closest — and the base paper we replicate")
note = find_shape(sn, "Minimum 8")
# Was: "author lists and page ranges sit behind IEEE Xplore ... not invented".
# No longer true -- every reference is resolved against the publisher record
# via Crossref, so the caveat is gone and this fits on one line.
set_body(note, [para([("[9] is the BASE PAPER: we reproduce its premise \u2014 a gate waveform "
                       "chosen by a multibit code \u2014 on GaN, then extend it. All five are "
                       "verified against the publisher record.", False)],
                     level=0, sz=1000, spc=0, bullet=False)])
tbl = table_of(sn)

# Column headings in the form the review panel asks for:
#   Ref | Author | Title & source | Aim (what they set out to do) |
#   Gap we found in it, and the part of that gap THIS project fills.
HEAD = ["Ref", "Author(s) & Year", "Title / Source",
        "Aim  —  what the paper sets out to do",
        "Gap we found  \u2192  what WE fill"]
for c, h in enumerate(HEAD):
    set_cell(tbl.cell(0, c), [(h, True)], size=Pt(9.5))

# What this project supplies that each paper leaves open. Kept to one
# sentence each: the panel reads the row, not an essay.
FILLS = {
 9:  "the exhaustive 720-word search that makes the fixed-vs-adaptive split "
     "measurable, and the GaN third-quadrant cost their SiC device never pays.",
 10: "the price of the pattern space: what an exhaustive search buys over their "
     "fixed pattern set.",
 11: "whether the setting has to change with the operating point at all. It "
     "does not: 3.9 % of baseline.",
 12: "both effects in ONE cost function, showing the objectives genuinely "
     "conflict (Pareto, slide 16).",
 13: "which field is worth scheduling. Only dead time, and only at light load; "
     "drive strength is worth 0.00 %.",
}

for row, ref in enumerate(CLOSEST, start=1):
    tag = "BASE" if ref.base else str(ref.n)
    set_cell(tbl.cell(row, 0), [(tag, ref.base)])
    set_cell(tbl.cell(row, 1), ref.table_author(), grey=not ref.done)
    set_cell(tbl.cell(row, 2), [(ref.title, False), ("\n" + ref.table_venue(), True)],
             size=Pt(9))
    set_cell(tbl.cell(row, 3), ref.method, size=Pt(9))
    gap = [(ref.finding + "  ", False)]
    if ref.n in FILLS:
        gap += [("WE FILL: ", True), (FILLS[ref.n], False)]
    set_cell(tbl.cell(row, 4), gap, size=Pt(9))
tbl._tbl.remove(tbl.rows[6]._tr)
for row in range(len(tbl.rows)):
    tbl.rows[row].height = Inches(0.88 if row else 0.35)
# 6.62 was measured against the table's DECLARED height (6.25 in). Cells wrap,
# so it actually renders to ~6.72 in and the footnote landed on the last row.
move_note(sn, 6.88)
# At 6.88 the full-width note runs under the page number in the bottom-right
# corner, so trim its width to stop short of it.
for _sh in sn.shapes:
    if _sh.has_text_frame and _sh.text_frame.text.startswith("[9] is the BASE PAPER"):
        _sh.width = Inches(11.30)
        break

# ---------------- slides 8 & 9 : results ----------------------------------
# Both arrive empty from the template (title + logo only). Content area runs
# from y=1.45 to y=6.85, x=0.70 to 12.60.

def add_text(slide, x, y, w, h, paras):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    bp = tb.text_frame._txBody.find(q("bodyPr"))
    for k in ("lIns", "tIns", "rIns", "bIns"):
        bp.set(k, "0")
    set_body(tb, paras)
    return tb

def stat(slide, x, y, w, big, label, sub):
    """A large number with a caption under it."""
    t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(1.30))
    tf = t.text_frame; tf.word_wrap = True
    bp = tf._txBody.find(q("bodyPr"))
    for k in ("lIns", "tIns", "rIns", "bIns"): bp.set(k, "0")
    set_body(t, [
        para([(big, True)], level=0, sz=3200, spc=100, bullet=False),
        para([(label, True)], level=0, sz=1300, spc=60, bullet=False),
        para([(sub, False)], level=0, sz=1150, spc=0, bullet=False),
    ])
    return t

# --- slide 8: the problem reproduced, and fixed ---------------------------
s8 = S[11]
s8.shapes.add_picture(RES + "/fig1_crosstalk.png",
                      Inches(0.70), Inches(1.45), width=Inches(7.30))
add_text(s8, 0.70, 6.30, 7.30, 0.45, [
    para([("Fig. 1  Spurious gate voltage on the off-device: the failure, and the actuator that "
           "removes it.", False)], level=0, sz=1100, spc=0, bullet=False)])

add_text(s8, 8.35, 1.45, 4.25, 0.45, [
    para([("The problem is real, and reproduced", True)], level=0, sz=1500, spc=0, bullet=False)])
stat(s8, 8.35, 2.00, 4.25, "1.65 V", "spurious gate peak vs 1.4 V threshold",
     "Fastest drive, no clamp — false turn-on")
stat(s8, 8.35, 3.35, 4.25, "2.58 V", "margin with clamp + −2 V off-bias",
     "The actuator removes it completely")
add_text(s8, 8.35, 4.75, 4.25, 2.00, [
    para([("And safety is nearly free", True)], level=0, sz=1500, spc=200, bullet=False),
    para([("The energy-optimal control word is already crosstalk-safe at three of four corners. "
           "At the fourth, safety costs ", False), ("0.04 % of switching energy", True),
          (" — 3.335 µJ against 3.334 µJ.", False)], level=0, sz=1300, spc=200, bullet=False),
    para([("The clamp itself is worth 9.7–12.2 % of the blended cost, from a static architecture "
           "choice needing no sensing or controller.", False)], level=0, sz=1300, spc=0, bullet=False),
])

# --- slide 9: the headline negative result -------------------------------
s9 = S[12]
s9.shapes.add_picture(RES + "/paper_fig2_ceiling.png",
                      Inches(0.70), Inches(1.55), width=Inches(6.60))
add_text(s9, 0.70, 5.35, 6.60, 0.45, [
    para([("Fig. 2  Cost of one fixed control word vs the true per-corner optimum, four corners.",
           False)], level=0, sz=1100, spc=0, bullet=False)])

add_text(s9, 0.70, 5.95, 6.60, 0.90, [
    para([("Full 720-word search at every corner — 2,880 transients — so each per-corner "
           "optimum is a true optimum, not the best of a shortlist.", False)],
         level=0, sz=1300, spc=0, bullet=False)])

add_text(s9, 7.65, 1.45, 4.95, 0.45, [
    para([("What adaptation is actually worth", True)], level=0, sz=1500, spc=0, bullet=False)])
stat(s9, 7.65, 2.00, 4.95, "5.2 %", "ceiling on operating-point scheduling",
     "vs the best single fixed word (= 3.9 % of baseline). A fixed word is nearly as good.")
# Grown from 3.40 to 3.90 in and narrowed 4.95 -> 4.80 to make room for the
# 36-point-grid disclosure; 4.80 keeps the right edge clear of the page number.
add_text(s9, 7.65, 3.32, 4.95, 3.72, [
    para([("The benefit is not spread out. ", True),
          ("Three corners lose 1–4 % from a fixed word; one loses 12.7 %.", False)],
         level=0, sz=1300, spc=180, bullet=False),
    para([("5.2 % is the generous figure: ", True),
          ("the denser 36-point grid gives ", False), ("2.0 %", True), (".", False)],
         level=0, sz=1200, spc=180, bullet=False),
    para([("It is carried by the dead time — and the dead time by one corner. ", True),
          ("Freezing dead time costs 5.45 % across four corners. Drop the light-load "
           "50 V / 2 A corner and it costs ", False), ("0.00 %", True),
          (": three corners want 5 ns, only that one wants 15 ns. So the adaptive hardware "
           "reduces to a light-load detector picking one of two dead times. Freezing pull-up "
           "drive strength costs ", False), ("0.00 %", True),
          (" — and drive strength is what the active-gate-driver literature actually schedules.",
           False)], level=0, sz=1200, spc=180, bullet=False),
    para([("Robust to the device, conditional on the layout. ", True),
          ("Across 21,600 transients, no device parameter moves the ceiling outside 4.3–7.7 %. "
           "Halving the loop inductance takes it to 13.5 %, and loop inductance is board layout, "
           "not the transistor [3].", False)], level=0, sz=1300, spc=0, bullet=False),
])
p.save(OUT)
print("slides 8 and 9 built")

# ---------------- references slide -----------------------------------------
ref_shape = find_shape(S[19], "A. Author")
# 13 refs no longer fit at 11.5 pt with 3 pt leading; tighten both, and mark
# the base paper so the panel can find it in the list
paras = []
for r in R.REFS:
    head = "[%d]  " % r.n
    runs = [(head, True)]
    if r.base:
        runs.append(("BASE PAPER — ", True))
    runs.append((r.cite(), False))
    if r.url():
        runs.append(("   " + r.url(), True))
    paras.append(para(runs, level=0, sz=900, spc=120, bullet=False))
n_done, n_pend = len(R.DONE), len(R.PENDING)
if n_pend == 0:
    note = [("All %d verified against the publisher record " % len(R.REFS), False),
            ("— authors, volume, issue and pages resolved through Crossref by DOI and title, "
             "not transcribed by hand.", False)]
else:
    note = [("%d of %d verified against the publisher record. " % (n_done, len(R.REFS)), False),
            ("The remaining %d carry title, journal status and Xplore document ID; " % n_pend, False),
            ("volume, issue, pages and authors were not reachable and are left blank rather "
             "than guessed.", False)]
paras.append(para(note, level=0, sz=850, spc=0, bullet=False))
set_body(ref_shape, paras)

# The template ships "Use IEEE format. Every reference listed must be cited in
# the slides / report." That is an instruction to the student, not something a
# panel should be shown, and cloning put it on both reference slides. Blank it.
try:
    n2 = find_shape(S[19], "Use IEEE format")
    set_body(n2, [para([("", False)], level=0, sz=1200, spc=0, bullet=False)])
except KeyError:
    pass

# ================= new slides 10 (demo) and 11 (evidence) ==================
# Both were cloned from the empty "Results (contd.)" slide, so each already
# carries the title box, the VIT logo and a page-number box — all of which
# still hold slide 9's values and must be corrected.

def retitle(slide, text):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == "Results (contd.)":
            for pa in sh.text_frame.paragraphs:
                for r in pa.runs:
                    r.text = text; return

def renumber(slide, n):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.left and sh.left > Inches(12.0) \
           and sh.top and sh.top > Inches(6.8):
            for pa in sh.text_frame.paragraphs:
                for r in pa.runs:
                    r.text = str(n); return

# ---- slide 10 : the demo video ------------------------------------------
s10 = S[16]
retitle(s10, "Demo")
# pipeline_demo.mp4 trimmed to 17.5 s. The original's closing caption read
# "the Miller clamp buys 14.7 %", the superseded 36-corner shortlist figure;
# the audit puts it at 9.7-12.2 %. The terminal body is untouched and every
# number in it still reproduces. research_demo.mp4 is not used at all: its
# closing caption says the benefit "collapses to one bit - the off-bias rail",
# which the freeze test corrected to dead time.
# The generator (scripts/make_demo_video.py) writes to results/. Naming the
# file bare resolved it against review/, where a stale copy sat: the video
# was regenerated and the deck silently kept embedding the old one. Take it
# from the same directory the generator writes to, so the two cannot drift.
VID = os.path.join(RES, "demo_crosstalk_explained.mp4")
# 1200x700 source -> 1.714 aspect. 7.60 in wide gives 4.43 in tall.
s10.shapes.add_movie(VID, Inches(0.70), Inches(1.55), Inches(7.60), Inches(4.43),
                     poster_frame_image="poster_crosstalk.png", mime_type="video/mp4")
add_text(s10, 0.70, 6.15, 7.60, 0.40, [
    para([("Click to play (22 s). Real ngspice waveforms from sim/dpt.cir, run twice "
           "— only the pacing and the captions are presentational.", False)],
         level=0, sz=1100, spc=0, bullet=False)])

add_text(s10, 8.65, 1.55, 3.95, 4.60, [
    para([("What it shows", True)], level=0, sz=1500, spc=240, bullet=False),
    para([("TOP row: the switch node. The low-side device turns on and V(sw) "
           "collapses from 100 V in a few nanoseconds. That slew is the cause.",
           False)], level=0, sz=1250, spc=170, bullet=True),
    para([("BOTTOM row: the OFF device's gate. The dV/dt couples through C", False),
          ("GD", False), (" and lifts it.", False)],
         level=0, sz=1250, spc=170, bullet=True),
    para([("LEFT, failing: the gate reaches ", False), ("+1.65 V", True),
          (" — above the 1.4 V threshold drawn on the plot. The device turns on when it "
           "must not.", False)], level=0, sz=1250, spc=170, bullet=True),
    para([("RIGHT, shipped: same circuit, same word, plus the Miller clamp and −2 V "
           "off-bias. Peak ", False), ("−1.18 V", True),
          (", a 2.58 V margin.", False)], level=0, sz=1250, spc=170, bullet=True),
    para([("Both panels are the same simulation deck; only CLKEN and VNEG differ.",
           False)], level=0, sz=1150, spc=0, bullet=False),
])

# ---- slide 11 : the evidence behind the numbers --------------------------
s11 = S[17]
retitle(s11, "Why the numbers hold")

TILES = [
    ("34,622", "transient simulations",
     "Every sweep, corner study, robustness run, loop-inductance and EMI sweep, start to finish."),
    ("720", "control words, searched in full",
     "At every corner — so each per-corner optimum is a true optimum, not the best of a shortlist."),
    ("25×", "timestep refinement survived",
     "Every reported quantity must be flat across it. Enforced, not assumed."),
    ("10", "wrong numbers caught, and corrected",
     "By that discipline, before any of them reached the report. Every one came from resampling "
     "or re-deriving a figure of our own — none from a reviewer."),
]
X0, Y0, W, DX, DY = 0.70, 1.70, 5.75, 6.15, 2.35
for i, (big, label, sub) in enumerate(TILES):
    x = X0 + (i % 2) * DX
    y = Y0 + (i // 2) * DY
    t = s11.shapes.add_textbox(Inches(x), Inches(y), Inches(W), Inches(2.00))
    bp = t.text_frame._txBody.find(q("bodyPr"))
    for k in ("lIns", "tIns", "rIns", "bIns"): bp.set(k, "0")
    t.text_frame.word_wrap = True
    set_body(t, [
        para([(big, True)], level=0, sz=4000, spc=120, bullet=False),
        para([(label, True)], level=0, sz=1500, spc=100, bullet=False),
        para([(sub, False)], level=0, sz=1250, spc=0, bullet=False),
    ])

add_text(s11, 0.70, 6.45, 11.90, 0.45, [
    para([("Stated plainly: ", True),
          ("LTspice has now been run: the three shipped .cir files reproduce 1.6487 / 0.8282 / "
           "\u22121.1768 V, matching ngspice to within 2 mV. Spectre is still a port \u2014 "
           "no installation available.", False)],
         level=0, sz=1150, spc=0, bullet=False)])

# ---- page numbers on every slide from 10 onward --------------------------
for n, sl in enumerate(S, start=1):
    if n >= 5:
        renumber(sl, n)

p.save(OUT)
print("slides 10 (demo, video embedded) and 11 (evidence) built; page numbers fixed")

# ---------------- geometry fixes found by qa.py ---------------------------
# The project title is longer than the template's stub, so it wraps to two
# lines at 32 pt and needs a taller box. There is 2.35 in of clear space below
# it, so growing it cannot collide with anything.
for sh in S[1].shapes:
    if sh.has_text_frame and sh.text_frame.text.startswith("A Segmented Gate Driver"):
        sh.height = Inches(1.30)

# The RTL slide's text box inherited the template's full 5.40 in height, so it
# ran under the waveform figure that now sits at y=4.05.
for sh in S[10].shapes:
    if sh.has_text_frame and sh.text_frame.text.startswith("Three modules emitting"):
        sh.height = Inches(2.45)

# The reference list inherited the template's full-height content box (5.40 in)
# but holds about 4 in of text, so the tail ran under the footnote.
for sh in S[19].shapes:
    if sh.has_text_frame and sh.text_frame.text.startswith("[1]"):
        sh.height = Inches(4.75)

p.save(OUT)
print("geometry fixes applied")


# ================= slide 11 : robustness ==================================
s_rob = S[15]
retitle(s_rob, "Backup — the device doesn\u2019t move it; layout does")

ROWS = [("Loop inductance −50 %  (1.5 nH)", "13.54 %", "+7.59", True),
        ("Input capacitance +30 %",              "7.71 %", "+1.76", False),
        ("Miller capacitance −50 %",        "6.93 %", "+0.99", False),
        ("nominal",                              "5.95 %",     "—", None),
        ("Threshold −20 %",                 "5.64 %", "−0.31", False),
        ("Transconductance +30 %",               "5.08 %", "−0.87", False),
        ("Threshold +20 %",                      "4.90 %", "−1.05", False),
        ("Input capacitance −30 %",         "4.76 %", "−1.19", False),
        ("Transconductance −30 %",          "4.45 %", "−1.49", False),
        ("Miller capacitance +50 %",             "4.32 %", "−1.63", False),
        ("Loop inductance +50 %  (4.5 nH)",      "0.55 %", "−5.39", True)]

# Three separate boxes per row. Space-padding does not align columns in a
# proportional font, so the numeric columns get their own x positions.
COL_L, COL_C, COL_D = 0.70, 4.55, 5.75
def row(y, label, ceil, delta, bold, sz=1200):
    add_text(s_rob, COL_L, y, 3.80, 0.30,
             [para([(label, bold)], level=0, sz=sz, spc=0, bullet=False)])
    add_text(s_rob, COL_C, y, 1.10, 0.30,
             [para([(ceil, bold)], level=0, sz=sz, spc=0, bullet=False)])
    add_text(s_rob, COL_D, y, 1.20, 0.30,
             [para([(delta, False)], level=0, sz=sz, spc=0, bullet=False)])

row(1.45, "Perturbation", "Ceiling", "vs nom.", True, 1150)
y = 1.88
for label, ceil, delta, big in ROWS:
    row(y, label, ceil, delta, big is True or big is None)
    y += 0.335

add_text(s_rob, 7.55, 1.45, 5.05, 0.60, [
    para([("Robust to the device,", True)], level=0, sz=1500, spc=0, bullet=False),
    para([("conditional on the layout.", True)], level=0, sz=1500, spc=0, bullet=False)])
add_text(s_rob, 7.55, 2.25, 5.05, 4.50, [
    para([("21,600 transients", True), (" — a full 720-word search at two corners under each "
           "of eleven single-parameter perturbations. One failure.", False)],
         level=0, sz=1300, spc=200, bullet=True),
    para([("Every ", False), ("device", True),
          (" parameter leaves the ceiling between 4.3 % and 7.7 % — at most 1.76 points from "
           "nominal. The objection this study was built to answer is not supported.", False)],
         level=0, sz=1300, spc=200, bullet=True),
    para([("The one thing that moves it is ", False), ("loop inductance", True),
          (", which is board layout, not the transistor. A tighter loop commutates faster, so "
           "more charge couples through C", False), ("GD", False),
          (": the median spurious gate voltage goes 0.640 V to 3.522 V and the feasible set "
           "collapses 474 → 158.", False)], level=0, sz=1300, spc=200, bullet=True),
    para([("So the claim is narrower and more useful than “robust”: from about 2.5 nH "
           "upward a fixed word is nearly as good; on a tighter loop the question has to be "
           "re-asked, and the answer there is not even monotonic — see slide 19.", False)],
         level=0, sz=1300, spc=200, bullet=True),
    para([("Baseline note. ", True),
          ("Nominal here is 5.95 %, not the 5.2 % headline — a two-corner search on its "
           "own sweep. Compare the ", False), ("vs nom.", True),
          (" column, not the absolute values.", False)],
         level=0, sz=1150, spc=0, bullet=True),
])

# ================= slide 14 : conclusion ==================================
s_con = S[18]
set_title(s_con, "Conclusion & next steps")
con_shape = find_shape(s_con, "Problem Statement:")
set_body(con_shape, [
    para([("What the data supports", True)], level=0, sz=1700, spc=240, bullet=False),
    para([("Choosing the word well: ", True), ("25.1 %.   ", False),
          ("Adapting it per operating point: ", True), ("3.9 %.", False)],
         level=0, sz=1600, spc=220, bullet=True),
    para([("One comparator takes 72 % of that 3.9 %.", True)],
         level=0, sz=1600, spc=220, bullet=True),
    para([("The deliverable: ", True),
          ("a fixed word plus a light-load comparator. The full sense + ADC + lookup table "
           "is left justifying ", False), ("3.7 %", True), (".", False)],
         level=0, sz=1600, spc=340, bullet=True),
    para([("What is next", True)], level=0, sz=1700, spc=240, bullet=False),
    para([("Review-II \u2014 ", True),
          ("transistor-level output stage in Cadence; re-run the ceiling on real devices.",
           False)], level=0, sz=1600, spc=220, bullet=True),
    para([("Review-III \u2014 ", True),
          ("measure a hardware half-bridge. Until then this is a simulation study, and is "
           "titled as one.", False)], level=0, sz=1600, spc=220, bullet=True),
    para([("Limits we state ourselves: ", True),
          ("no silicon measured; one device model underlies everything.", False)],
         level=0, sz=1450, spc=0, bullet=True),
])

for n, sl in enumerate(S, start=1):
    renumber(sl, n)
p.save(OUT)
print("slides 11 (robustness) and 14 (conclusion) built")

# ---------------- slide 2 : name / guide fields ---------------------------
# The label and its value live in SEPARATE shapes, so a run-follows-run search
# within one shape never finds them. Kept last so nothing can clobber it.
REPL = {
    "Name — Reg. No.":        "Sanjay Kumar 23BEC1447  ·  Aamir Abdullah 23BPS1197  ·  Amritha S 23BEC1368",
    "Dr. Guide Name, School": "Dr. Bindu  —  SENSE",
}
hits = 0
for sh in S[1].shapes:
    if not sh.has_text_frame: continue
    for pa in sh.text_frame.paragraphs:
        for r in pa.runs:
            if r.text.strip() in REPL:
                r.text = REPL[r.text.strip()]; hits += 1

# Three names and reg numbers do not fit the template's one-line value box
# (4.40 in wide, 0.42 in tall, sized for a single student). Widen it to the
# signature block's left edge and give it the second line it needs; the row
# below starts at 4.00 in, so 0.55 in still clears it.
# The template draws a fixed underline rule under each value field, so a
# second line is struck through by it. Keep the names to ONE line: widen the
# box to the signature block's edge and drop this field to 9.5 pt. qa.py does
# not catch that collision -- the rule is a separate shape, not text.
for sh in S[1].shapes:
    if sh.has_text_frame and "23BEC1447" in sh.text_frame.text:
        sh.width = Inches(5.05)
        for pa in sh.text_frame.paragraphs:
            for r in pa.runs:
                r.font.size = Pt(9.5)
        break

# ---------------- informative titles on the two results slides ------------
for sl, txt in ((S[11], "Result 1 — crosstalk is real; the clamp fixes it"),
                (S[12], "Result 2 — scheduling is worth only 5.2 %")):
    for sh in sl.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() in ("Results", "Results (contd.)"):
            for pa in sh.text_frame.paragraphs:
                for r in pa.runs:
                    r.text = txt
            break

p.save(OUT)
print("slide 2: %d of %d fields set; results slides retitled" % (hits, len(REPL)))

# ================= speaker notes + time budget ============================
# 17 slides against an 8-10 minute slot. Without a budget the talk overruns
# and the marker docks presentation quality, so each slide carries its target
# and the two droppable slides say so.
# Official brief: 10 minutes, and the mark split is problem clarity 2,
# literature 1, proposed solution 1, presentation quality 1.
#
# The CORE set below totals 535 s = 8.9 min, leaving ~1 min for transitions
# and the inevitable overrun. Three slides are marked BACKUP: they stay in the
# deck and are strong answers to likely questions, but they are not presented
# unless asked. Presentation quality is a mark, and overrunning a 10-minute
# slot is the easiest way to lose it.
# 20 slides, 10-minute slot. CORE totals 570 s = 9.5 min, leaving 30 s of
# margin. Four slides are BACKUP: in the deck, prepared, not presented.
NOTES = {
 2:  ("15 s", "CORE. Project name and the one-line claim: a segmented GaN gate driver, and a "
               "measurement of what per-operating-point adaptation is actually worth."),
 3:  ("skip", "SKIP. Template slide."),
 4:  ("60 s", "CORE — the 2-mark slide, the biggest block in the review. Land two things: the "
              "failure is real and specific (1.65 V on a 1.40 V threshold), and the gap is "
              "that prior work conflates 'better fixed settings' with 'settings that adapt'."),
 5:  ("30 s", "CORE — literature mark. Do not read the table. Eight references, three verified "
              "to page level; point at [3], layout sets loop inductance, which decides our "
              "answer."),
 6:  ("25 s", "CORE — literature mark. Closest prior work. [8] adapts dead time, the one field "
              "we find worth scheduling — and it does so across a 0.2-2 A range, which is "
              "exactly the light-load corner our leave-one-out test shows carries all of that "
              "benefit. Independent support, not competition."),
 7:  ("30 s", "CORE — proposed-solution mark. The 720-point control word and the method: "
              "exhaustive search at every corner, so each per-corner optimum is a TRUE optimum "
              "and not the best of a shortlist."),
 8:  ("30 s", "CORE. Walk the diagram left to right in one sentence, then stop on the dashed "
              "block: that is the sensing and lookup-table machinery the adaptive premise "
              "needs. The whole project is a measurement of what it buys. Let the picture do "
              "the work — do not narrate every box."),
 9:  ("40 s", "CORE — the 50 %-completion requirement. Five claims, each with a number. If "
              "asked 'is this really 50 %?': the search is complete, stress-tested, and both "
              "objections to it - layout and EMI - have been tested rather than argued. What "
              "remains is silicon and hardware."),
 10: ("30 s", "CORE — the timeline and tools requirement. M1-M6. If asked what could go wrong: "
              "M1, the PDK — a 1.8 V / 3.3 V teaching PDK cannot take a 5 V rail."),
 11: ("BACKUP", "BACKUP — the FPGA. If asked: written AND verified, eight asserted properties, "
                "then mutation-tested. Dead time gets a live register and drive strength is "
                "strapped, because that is what the study found."),
 12: ("40 s", "CORE. The problem reproduced and fixed. Point at the red X below the threshold "
              "line. Then: safety costs 0.04 % — nearly free once the clamp is present."),
 13: ("25 s", "CORE. 5.2 %. Say plainly this is a NEGATIVE result about the adaptive premise "
              "and you are reporting it rather than hiding it. Keep it short — slides 18 and "
              "19 do the real work."),
 14: ("45 s", "CORE — the novelty slide, and the one sentence the panel should leave with: "
              "'a full sense-plus-ADC-plus-lookup-table is left justifying 3.7 % of the gain "
              "over a fixed word and one comparator.' Build it in three steps. (1) Choosing "
              "the fixed word well is worth 25.1 %. (2) Adapting per corner adds only 3.9 % "
              "— a seventh of the total. (3) And 72 % of even THAT is capturable with a "
              "single light-load comparator, not a lookup table. No published work separates "
              "these, because separating them needs the exhaustive search rather than a "
              "shortlist. If asked why nobody found this: they report one number."),
 15: ("40 s", "CORE — the design chart, and the most USEFUL thing in the deck. Do NOT call it "
              "a single crossover: eight points show a BAND. Adaptive control pays from about "
              "2.5 nH down, peaks at 13.5 % at 1.5 nH, then FALLS BACK to 8.1 % at 1.0 nH. Say "
              "why, because it is the interesting part: at 1.0 nH only 165 of 720 words are "
              "still safe, so the fixed word and the per-corner optima are squeezed into the "
              "same narrow region and the gap scheduling exploits shrinks. Below ~2 nH "
              "feasibility binds, not optimisation. We checked it is not the clamp-chatter "
              "artefact — excluding every chatter point leaves all eight ceilings unchanged to "
              "the digit. The literature does not state this trade-off because it never "
              "separates the fixed and adaptive halves."),
 16: ("BACKUP", "BACKUP — the robustness table. Best answer to 'is this just your model?': "
                "every device parameter leaves the ceiling at 4.3-7.7 %; only loop inductance "
                "moves it, and that is board layout, not the transistor."),
 17: ("25 s", "CORE. Play the video. Real captured ngspice output — the failure, the fix, the "
              "search. Let it run; do not talk over it."),
 18: ("BACKUP", "BACKUP — use if the panel probes rigour. TEN wrong numbers caught by our own "
                "checks is a strength; say it that way. The best one to tell: dead time ranked "
                "first at 5.45 % on four corners and dead LAST at 0.00 % on two, and we chased "
                "it to a single light-load corner rather than quoting the flattering number. The "
                "leave-one-out table is in results/FINDINGS.md section 32 if they want it."),
 19: ("35 s", "CORE. Close on what the data supports and what it now answers. Do NOT concede "
              "EMI as a limitation any more - it was tested: pricing it makes scheduling worth "
              "LESS, 5.95 % down to 0.15 %. State the remaining limits yourself instead: the "
              "two EMI measures disagree about drive strength, no silicon has been measured, "
              "and one device model underlies everything."),
 20: ("skip", "Reference list. Leave up during questions."),
 21: ("—",    "Thank you. Expect: 'why is the clamp worth more than scheduling?', 'have you run "
              "real silicon?' (no — Review-II), 'did Cadence actually run?' (no, and the deck "
              "says so), 'is 8 references enough?' (it is the stated minimum). Newer ones: "
              "'your dead-time number depends on which corners you pick' — yes, and we found "
              "that ourselves; it is one light-load corner - leave-one-out table in "
              "FINDINGS.md section 32. "
              "'Doesn't EMI change your answer?' — we ran it; it makes scheduling worth less, "
              "not more."),
}

NOTES_BODY = (
 '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
 '<p:nvSpPr><p:cNvPr id="99" name="Notes Placeholder"/>'
 '<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
 '<p:nvPr><p:ph type="body" idx="1"/></p:nvPr></p:nvSpPr>'
 '<p:spPr/><p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody></p:sp>')

def set_note(slide, text):
    """Cloned slides get a notes part with no body placeholder, so
    notes_text_frame is None until one is injected."""
    ns = slide.notes_slide
    if ns.notes_text_frame is None:
        ns.shapes._spTree.append(etree.fromstring(NOTES_BODY))
    ns.notes_text_frame.text = text


for n, (budget, note) in NOTES.items():
    ns = S[n - 1].notes_slide
    if ns.notes_text_frame is None:
        # Slides cloned by add_slide.py inherit a notes part with no body
        # placeholder, so python-pptx has nothing to write into. Add one.
        ns.shapes._spTree.append(etree.fromstring(NOTES_BODY))
    ns.notes_text_frame.text = "[%s]  %s" % (budget, note)

p.save(OUT)
print("speaker notes added to %d slides (budget: ~8 min of talking)" % len(NOTES))

# 2480x900 -> aspect 2.756. At 7.30 in wide the figure is 2.65 in tall.
S[10].shapes.add_picture(RES + "/fig_rtl_waveform.png",
                         Inches(3.00), Inches(4.05), width=Inches(7.30))
add_text(S[10], 3.00, 6.80, 7.30, 0.34, [
    para([("Icarus Verilog VCD. ls_pu falls to 0 before hs_pu rises — the shaded dead time is "
           "the gap, and no shoot-through window exists.", False)],
         level=0, sz=1000, spc=0, bullet=False)])

# ================= slide 8 : architecture =================================
s_arch = S[7]
retitle(s_arch, "System Architecture")
# Size by HEIGHT, not width. Sizing by width hard-codes an assumption about
# the figure's aspect, and when arch_diagram.py was redrawn (2.095 -> 1.969)
# the picture silently grew to 5.59 in tall and sat on its own caption.
# The band from 1.32 in to the caption at 6.66 in is 5.24 in.
s_arch.shapes.add_picture(RES + "/fig_architecture.png",
                          Inches(1.55), Inches(1.32), height=Inches(5.20))
# No slide-level caption here: the figure carries its own footer saying the
# same thing, and printing it twice is exactly the padding this deck is being
# stripped of.

# ================= slide 13 : the novelty =================================
s_nov = S[13]
retitle(s_nov, "Result 3 — the two halves nobody separates")

# The base paper is not just cited, it is IMPLEMENTED and run on the same
# testbench. models/basedrv.lib + scripts/basepaper_compare.py.
add_text(s_nov, 0.70, 1.38, 12.10, 0.80, [
    para([("The base paper [9] ", False), ("(doi.org/10.1002/cta.3136)", True),
          (" shows a gate waveform can be chosen by a digital code. It, and every "
           "active-gate-driver paper after it, then reports ", False), ("one", True),
          (" number — the improvement over a conventional driver. That number bundles two "
           "separable effects, and only the second needs sensing, an ADC and a lookup table. "
           "Separating them needs the exhaustive search, which is why nobody has.", False)],
         level=0, sz=1250, spc=0, bullet=False)])

NOV = [("(A)   Choose a better FIXED word",        "25.1 %", "no sensing · no LUT"),
       ("(B)   ADAPT it per operating point",      "3.9 %",  "needs all of it"),
       ("(B′)  …but ONE comparator captures 72 % of (B)", "2.8 %", "a threshold, not a LUT")]
for i, (label, val, sub) in enumerate(NOV):
    y = 2.42 + i * 1.12
    add_text(s_nov, 0.70, y, 6.55, 0.40,
             [para([(label, True)], level=0, sz=1400, spc=0, bullet=False)])
    add_text(s_nov, 7.45, y - 0.20, 2.15, 0.72,
             [para([(val, True)], level=0, sz=2600, spc=0, bullet=False)])
    add_text(s_nov, 9.85, y + 0.04, 2.95, 0.60,
             [para([(sub, False)], level=0, sz=1100, spc=0, bullet=False)])

add_text(s_nov, 0.70, 5.92, 12.10, 1.10, [
    para([("So the full sense + ADC + lookup table is left justifying ", False),
          ("3.7 % of the total achievable gain", True),
          (" over a fixed word plus one comparator. Adaptation is 13.4 % of the gain; a single "
           "threshold takes 72 % of that.", False)], level=0, sz=1300, spc=130, bullet=False),
    para([("Every figure shares one baseline — the median control word that is safe at all four "
           "corners. The split does not rest on the weighting: swept over 106 overshoot "
           "weights, (A) stays ", False), ("23.4–29.0 %", True), (" and (B) ", False),
          ("1.3–6.4 %", True), (" — and (A) exceeds (B) at ", False), ("every", True),
          (" weight tested, out to 5.0. ", False),
          ("scripts/novelty.py", True), (", ", False), ("scripts/weight_sensitivity.py", True),
          (".", False)], level=0, sz=1200, spc=0, bullet=False),
])

for n, sl in enumerate(S, start=1):
    renumber(sl, n)
p.save(OUT)
# ================= slide 15 : the design chart ============================
s_chart = S[14]
retitle(s_chart, "Result 4 — adaptive pays below ~2.5 nH")
# aspect 1.856 (two stacked panels); 1.35 in to the caption at 6.36 in is
# 4.90 in tall, so the width that fits is 4.90 * 1.856 = 9.09 in, centred.
s_chart.shapes.add_picture(RES + "/fig_lloop.png",
                           Inches(2.12), Inches(1.35), width=Inches(9.09))
add_text(s_chart, 0.70, 6.36, 11.90, 0.70, [
    para([("Eight inductances, 7,200 transients. The answer is a ", False), ("band", True),
          (", not a threshold: adaptive control pays only from about 2.5 nH down, peaking at "
           "13.5 % at 1.5 nH, and it falls back to 8.1 % at 1.0 nH because only 165 of 720 "
           "words stay safe there — below ~2 nH feasibility binds, not optimisation. ", False),
          ("A designer measures their loop and reads off the decision", True),
          (" — a trade-off the literature does not state.", False)],
         level=0, sz=1200, spc=0, bullet=False)])

for n, sl in enumerate(S, start=1):
    renumber(sl, n)
p.save(OUT)
print("slides 8 (architecture), 14 (novelty) and 15 (design chart) built")


# ================= slide 9 : what a control word IS =======================
# The deck's central object was never defined anywhere. Slides 13-15 are all
# about "the fixed word" versus "scheduling the word", and a reader had no
# way to know a word is a 6-tuple. Without this slide those results are
# unreadable, which is exactly how they read.
import copy as _copy

def clone_after(prs, src_idx, dest_idx):
    """Duplicate a slide and move the copy to dest_idx.

    Deep-copying the shape XML alone leaves r:embed ids pointing at
    relationships the new slide part does not have, so the logo vanishes;
    the relationships are re-created and the ids remapped below."""
    src = prs.slides[src_idx]
    new = prs.slides.add_slide(src.slide_layout)
    for shp in list(new.shapes):
        shp._element.getparent().remove(shp._element)
    NOTES_RT = ("http://schemas.openxmlformats.org/officeDocument/2006/"
                "relationships/notesSlide")
    idmap = {}
    for rid, rel in src.part.rels.items():
        # NEVER carry the notes relationship across: the clone would SHARE the
        # source slide's notes part, so writing the copy's speaker note
        # silently overwrites the original's.
        if rel.reltype == NOTES_RT:
            continue
        if rel.is_external:
            idmap[rid] = new.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            idmap[rid] = new.part.relate_to(rel.target_part, rel.reltype)
    R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    for shp in src.shapes:
        el = _copy.deepcopy(shp._element)
        for node in el.iter():
            for attr in ("embed", "link", "id"):
                key = R_NS + attr
                if key in node.attrib and node.attrib[key] in idmap:
                    node.attrib[key] = idmap[node.attrib[key]]
        new.shapes._spTree.append(el)
    lst = prs.slides._sldIdLst
    items = list(lst)
    lst.remove(items[-1])
    lst.insert(dest_idx, items[-1])
    return prs.slides[dest_idx]

s_def = clone_after(p, 11, 8)          # copy the blank "Results" frame -> slide 9
for sh in list(s_def.shapes):          # strip the architecture content off the copy
    if sh.has_text_frame and sh.text_frame.text.strip() and not sh.has_table:
        if sh.width > Inches(4) and sh.top > Inches(1.0):
            sh._element.getparent().remove(sh._element)
    elif sh.shape_type == 13 and sh.left < Inches(11):
        sh._element.getparent().remove(sh._element)
# retitle() only matches the template's literal "Results (contd.)", and the
# clone's title was already rewritten, so set it directly
for sh in s_def.shapes:
    if sh.has_text_frame and sh.top is not None and sh.top < Inches(1.0) \
       and sh.left is not None and sh.left < Inches(1.0):
        tf = sh.text_frame
        for extra in list(tf.paragraphs)[1:]:
            extra._p.getparent().remove(extra._p)
        pa = tf.paragraphs[0]
        runs = list(pa.runs)
        if runs:
            runs[0].text = "What is a \u201ccontrol word\u201d?"
            for r in runs[1:]:
                r._r.getparent().remove(r._r)
        break

add_text(s_def, 0.70, 1.25, 12.0, 0.45, [
    para([("Everything in this deck is measured over control words, so this is the "
           "definition the rest of the results rest on.", False)],
         level=0, sz=1300, spc=0, bullet=False)])

FIELDS = [("N_PU", "pull-up slices ON", "1 – 8", "how hard the device is turned ON"),
          ("N_PD,LS", "low-side pull-down slices", "1 – 8", "how hard it is turned OFF"),
          ("N_PD,HS", "high-side pull-down slices", "1, 4, 8", "the off device's grip on its gate"),
          ("t_dead", "dead time", "5 / 10 / 15 / 25 ns", "gap before the other device turns on"),
          ("CLK_EN", "active Miller clamp", "off / on", "shorts the gate during the other edge"),
          ("V_neg", "gate off-bias", "0 V / −2 V", "how far below threshold the gate is held")]
YT, DY = 2.05, 0.52
add_text(s_def, 0.70, YT - 0.34, 12.0, 0.30, [
    para([("ONE control word = one setting of these six fields", True)],
         level=0, sz=1250, spc=0, bullet=False)])
for i, (sym, name, vals, why) in enumerate(FIELDS):
    y = YT + i * DY
    add_text(s_def, 0.85, y, 1.60, 0.34,
             [para([(sym, True)], level=0, sz=1150, spc=0, bullet=False)])
    add_text(s_def, 2.50, y, 3.05, 0.34,
             [para([(name, False)], level=0, sz=1150, spc=0, bullet=False)])
    add_text(s_def, 5.60, y, 2.05, 0.34,
             [para([(vals, True)], level=0, sz=1150, spc=0, bullet=False)])
    add_text(s_def, 7.75, y, 4.85, 0.34,
             [para([(why, False)], level=0, sz=1150, spc=0, bullet=False)])

add_text(s_def, 0.70, 5.28, 12.0, 1.30, [
    para([("8 × 8 × 3 × 4 × 2 × 2  =  ", False), ("720 control words", True),
          (".  We simulate every one of them, at every corner — that is what makes this an "
           "exhaustive search rather than a shortlist.", False)],
         level=0, sz=1300, spc=140, bullet=False),
    para([("A FIXED word", True),
          (" is one word chosen once and strapped at power-up — no sensing. ", False),
          ("SCHEDULING", True),
          (" means sensing the operating point and switching to a different word for it, which "
           "is what needs an ADC and a lookup table. The whole project measures the gap between "
           "those two.", False)], level=0, sz=1300, spc=0, bullet=False)])


# the definition slide is created after NOTES is applied, so it gets its own
DEF_NOTE = ("[30 s]  CORE — do NOT skip this and do not rush it. Every number after "
            "this slide is a comparison between control words, and if the panel has not "
            "got the definition, slides 13 to 16 are noise to them. Read the six fields "
            "out loud, land 720, then land the one distinction that matters: a FIXED word "
            "is strapped at power-up and costs nothing; SCHEDULING needs sensing, an ADC "
            "and a lookup table. The whole project measures the gap between those two.")
if s_def.has_notes_slide or True:
    set_note(s_def, DEF_NOTE)

for n, sl in enumerate(S, start=1):
    renumber(sl, n)
p.save(OUT)
print("slide 9 (control-word definition) inserted; deck now %d slides" % len(S))

# ============ slide 10 : base paper vs this work ==========================
# A review wants the delta stated, not inferred. Every row is either a fact
# about the two designs or a number this project measured; nothing here
# claims we beat the base paper on ITS metric, because we have not run it.
s_cmp = clone_after(p, 11, 9)
for sh in list(s_cmp.shapes):
    if sh.has_text_frame and sh.text_frame.text.strip() and not sh.has_table:
        if sh.width > Inches(4) and sh.top > Inches(1.0):
            sh._element.getparent().remove(sh._element)
    elif sh.shape_type == 13 and sh.left < Inches(11):
        sh._element.getparent().remove(sh._element)
for sh in s_cmp.shapes:
    if sh.has_text_frame and sh.top is not None and sh.top < Inches(1.0) \
       and sh.left is not None and sh.left < Inches(1.0):
        tf = sh.text_frame
        for extra in list(tf.paragraphs)[1:]:
            extra._p.getparent().remove(extra._p)
        runs = list(tf.paragraphs[0].runs)
        if runs:
            runs[0].text = "Base paper vs this work — what we improve"
            for r in runs[1:]:
                r._r.getparent().remove(r._r)
        break

add_text(s_cmp, 0.70, 1.22, 12.0, 0.28, [
    para([("Base paper [9]: Takayama & Hikihara, Int. J. Circuit Theory Appl., 50(1):183–196, "
           "2022 — a gate waveform chosen by a multibit digital code.", False)],
         level=0, sz=1150, spc=0, bullet=False)])

CMP = [
 ("Device",            "SiC MOSFET",                     "GaN HEMT (e-mode)  —  no body diode, so off-bias costs third-quadrant drop"),
 ("Control handle",    "multibit gate code, one field",  "6-field, 720-point control word  —  strength, dead time, clamp, off-bias"),
 ("Implementation",    "digital driver",                 "FPGA RTL in Verilog  ·  8 asserted properties, mutation-tested"),
 ("How codes chosen",  "selected codes, measured",       "EXHAUSTIVE  —  all 720 words at every corner, 34,622 transients"),
 ("Crosstalk",         "not addressed",                  "1.65 V spurious gate  →  −1.18 V  (2.58 V margin), clamp + off-bias"),
 ("What is reported",  "one number vs a conventional driver",
                       "the number SPLIT: fixed word 25.1 %, adapting it 3.9 %, one comparator takes 72 % of that"),
 ("Design rule",       "none",                           "adaptive control pays only below ~2.5 nH loop inductance"),
]
X0, X1, X2 = 0.70, 3.45, 7.05
YT, DY = 1.95, 0.60
for lbl, xs, ys in ((("Parameter", True), X0, YT - 0.42),):
    pass
add_text(s_cmp, X0, YT - 0.46, 2.60, 0.30,
         [para([("Parameter", True)], level=0, sz=1150, spc=0, bullet=False)])
add_text(s_cmp, X1, YT - 0.46, 3.45, 0.30,
         [para([("Base paper [9]", True)], level=0, sz=1150, spc=0, bullet=False)])
add_text(s_cmp, X2, YT - 0.46, 5.55, 0.30,
         [para([("This work", True)], level=0, sz=1150, spc=0, bullet=False)])
for i, (par, base, ours) in enumerate(CMP):
    y = YT + i * DY
    add_text(s_cmp, X0, y, 2.65, 0.50,
             [para([(par, True)], level=0, sz=1050, spc=0, bullet=False)])
    add_text(s_cmp, X1, y, 3.50, 0.50,
             [para([(base, False)], level=0, sz=1050, spc=0, bullet=False)])
    add_text(s_cmp, X2, y, 5.60, 0.50,
             [para([(ours, False)], level=0, sz=1050, spc=0, bullet=False)])

add_text(s_cmp, 0.70, 6.28, 12.0, 0.75, [
    para([("The parameter we improve is not a volt or a nanosecond — it is ", False),
          ("what the designer knows before committing silicon", True),
          (". [9] shows a digital code works. We measure what the code is worth, and how much "
           "of that worth needs sensing at all: 3.7 % of the achievable gain, which is what the "
           "sense + ADC + lookup table has to justify.", False)],
         level=0, sz=1200, spc=0, bullet=False)])

set_note(s_cmp, (
    "[35 s]  CORE — the comparison the panel will ask for. Land three rows and move on: "
    "we go from ONE field to six (720 words); from selected codes to an EXHAUSTIVE search; "
    "and from one bundled number to the number SPLIT into fixed versus adaptive. Be honest "
    "if pushed: we have NOT re-run [9]'s SiC experiment, so this is a comparison of method "
    "and scope, not a claim that we beat their result on their metric. That replication is "
    "the next step."))

for n, sl in enumerate(S, start=1):
    renumber(sl, n)
p.save(OUT)
print("slide 10 (base paper vs this work) inserted; deck now %d slides" % len(S))


# ============ slide 9 : the circuit itself =================================
# The architecture slide shows what talks to what. A panel asking "what did
# you actually simulate?" wants the schematic, so it goes right after it and
# before the control-word definition.
s_ckt = clone_after(p, 11, 8)
for sh in list(s_ckt.shapes):
    if sh.has_text_frame and sh.text_frame.text.strip() and not sh.has_table:
        if sh.width > Inches(4) and sh.top > Inches(1.0):
            sh._element.getparent().remove(sh._element)
    elif sh.shape_type == 13 and sh.left < Inches(11):
        sh._element.getparent().remove(sh._element)
for sh in s_ckt.shapes:
    if sh.has_text_frame and sh.top is not None and sh.top < Inches(1.0) \
       and sh.left is not None and sh.left < Inches(1.0):
        tf = sh.text_frame
        for extra in list(tf.paragraphs)[1:]:
            extra._p.getparent().remove(extra._p)
        runs = list(tf.paragraphs[0].runs)
        if runs:
            runs[0].text = "The circuit that is simulated"
            for r in runs[1:]:
                r._r.getparent().remove(r._r)
        break

s_ckt.shapes.add_picture(RES + "/fig_circuit.png",
                         Inches(2.22), Inches(1.22), width=Inches(8.89))
add_text(s_ckt, 0.70, 6.20, 12.10, 0.85, [
    para([("Every element is in ", False), ("sim/dpt.cir", True),
          (", which is the file to open in LTspice — it carries the active Miller clamp "
           "(S", False), ("clk", False), (" through 0.5 \u03a9). The three ", False),
          ("ltspice/*.asc", True),
          (" sheets illustrate the crosstalk mechanism only and do not draw the clamp", False),
          (". C", False), ("GD", False),
          (" on Q2 in red is the crosstalk path; GaN has no body diode, so holding the gate "
           "below threshold is paid for again across the dead time.", False)],
         level=0, sz=1200, spc=0, bullet=False)])
set_note(s_ckt, "[30 s]  CORE — use this the moment anyone asks what was actually simulated. "
                "Trace the red path with a finger: Q1 switches, the SW node moves, charge "
                "couples through C_GD into Q2's gate, and Q2 turns on when it should be off. "
                "Then point at the clamp and the off-bias mux — those are the two things that "
                "fix it. Do not read the component values out.")

for n, sl in enumerate(S, start=1):
    renumber(sl, n)
p.save(OUT)
print("slide 9 (circuit schematic) inserted; deck now %d slides" % len(S))


# ============ slide : MATLAB results ======================================
# The MATLAB work existed as scripts and loose PNGs and appeared nowhere in
# the deck. Both figures were regenerated by running the scripts.
s_mat = clone_after(p, 11, 15)
for sh in list(s_mat.shapes):
    if sh.has_text_frame and sh.text_frame.text.strip() and not sh.has_table:
        if sh.width > Inches(4) and sh.top > Inches(1.0):
            sh._element.getparent().remove(sh._element)
    elif sh.shape_type == 13 and sh.left < Inches(11):
        sh._element.getparent().remove(sh._element)
for sh in s_mat.shapes:
    if sh.has_text_frame and sh.top is not None and sh.top < Inches(1.0) \
       and sh.left is not None and sh.left < Inches(1.0):
        tf = sh.text_frame
        for extra in list(tf.paragraphs)[1:]:
            extra._p.getparent().remove(extra._p)
        runs = list(tf.paragraphs[0].runs)
        if runs:
            runs[0].text = "MATLAB — the Pareto front and the model check"
            for r in runs[1:]:
                r._r.getparent().remove(r._r)
        break

# MATLAB Online's own renders, from results/matlab_online/. They come out on a
# dark theme; a dark backing panel behind each makes that read as deliberate
# rather than as a figure that clashes with the slide.
# Both figures are MATLAB Online's own renders (results/matlab_online/) and are
# 4:3. The panel is sized to that aspect so the dark theme reads as deliberate
# instead of leaving black margins. MATLAB's crosstalk_model render is NOT used:
# it came out 219x605, a degenerate sliver, so the weight-sensitivity figure --
# a result the deck did not previously show -- takes the right-hand slot.
_FIGH = 3.60
_FIGW = _FIGH * 4.0 / 3.0
for _x, _png in ((0.70, "pareto_matlab.png"), (7.05, "fig_master_weight.png")):
    _bg = s_mat.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(_x - 0.10), Inches(1.20),
                                 Inches(_FIGW + 0.20), Inches(_FIGH + 0.20))
    _bg.fill.solid(); _bg.fill.fore_color.rgb = RGBColor(0x0B, 0x0B, 0x0B)
    _bg.line.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)
    _bg.shadow.inherit = False
    if _bg.has_text_frame:
        _bg.text_frame.text = ""
    s_mat.shapes.add_picture(RES + "/matlab_online/" + _png,
                             Inches(_x), Inches(1.30), height=Inches(_FIGH))
add_text(s_mat, 0.60, 5.10, 5.95, 1.05, [
    para([("720 words, 504 feasible (70 %).", True),
          (" A 7-word Pareto front; buying one point of overshoot back costs "
           "0.039 µJ. The green marker is the word the cost function picks.", False)],
         level=0, sz=1150, spc=0, bullet=False)])
add_text(s_mat, 6.95, 5.10, 5.75, 1.05, [
    para([("(A) stays above (B) at every weight tested.", True),
          (" Choosing the word well is worth 22.5–29.0 % of baseline across the whole "
           "range; adapting it per operating point 1.4–12.7 %. The ordering never "
           "flips.", False)], level=0, sz=1150, spc=0, bullet=False)])
add_text(s_mat, 0.60, 6.28, 12.10, 0.75, [
    para([("Dead-time margin saturates at about 15 ns", True),
          (": the marginal gain runs 269 → 31.5 → 1.8 → 0.0 mV/ns across 10/15/25/35 ns. "
           "Past 15 ns dead time costs conduction loss and buys nothing. ", False),
          ("results/gan_analysis.m", True), (", ", False),
          ("results/gan_master.m", True), (" — both run in MATLAB or Octave.",
           False)], level=0, sz=1200, spc=0, bullet=False)])
set_note(s_mat, "[35 s]  CORE — this is the MATLAB evidence, and it answers 'did you only "
                "use SPICE?'. Land two things. Left: 720 words searched, 504 feasible, and "
                "the Pareto front is only 7 words wide — the choice really is that "
                "constrained. Right: two textbook hand calculations for the same spurious "
                "voltage disagree by 8.7x, and the SPICE number sits between them. That is "
                "the argument for simulating. If you have time, add the dead-time "
                "saturation at 15 ns from the caption.")

# ---------------- FLOW CHART: one use case, start to finish ----------------
# Asked for explicitly. The architecture diagram says what talks to what and
# the circuit says what is wired to what; neither says what HAPPENS, in order,
# on one switching edge. This slide is that, for a single named use case.
_fl_idx = None
for _i, _sl in enumerate(S):
    for _sh in _sl.shapes:
        if _sh.has_text_frame and _sh.text_frame.text.strip().startswith("System Architecture"):
            _fl_idx = _i; break
    if _fl_idx is not None:
        break
if _fl_idx is None:
    raise SystemExit("architecture slide not found for flow-chart insertion")

s_flow = clone_after(p, _fl_idx, _fl_idx)          # sits BEFORE the architecture
set_title(s_flow, "How it works \u2014 one use case, start to finish")
for _sh in list(s_flow.shapes):
    if _sh.has_text_frame:
        _t = _sh.text_frame.text.strip()
        if _t and not _t.isdigit() and not _t.startswith("How it works"):
            _sh._element.getparent().remove(_sh._element)
    elif _sh.shape_type is not None and not _sh.has_text_frame and _sh.name != "Image 0":
        _sh._element.getparent().remove(_sh._element)
# Size by HEIGHT, not width: at 12.45 in wide this 1.90:1 figure is 6.3 in
# tall and runs off the bottom of a 7.5 in slide, over the caption and the
# page number.
s_flow.shapes.add_picture(RES + "/fig_flow.png",
                          Inches(1.86), Inches(1.30), height=Inches(5.05))
add_text(s_flow, 0.70, 6.55, 11.20, 0.45, [
    para([("Read the top row as the controller and the bottom row as the circuit. ", True),
          ("The one shaded diamond is the entire adaptive machinery.", False)],
         level=0, sz=1200, spc=0, bullet=False)])
for n, sl in enumerate(S, start=1):
    renumber(sl, n)
p.save(OUT)
print("flow-chart slide inserted; deck now %d slides" % len(S))

# ---------------- TECH STACK: what each tool actually produced -------------
# Asked for explicitly: show the software, and what each one contributed. The
# point of the slide is that no single tool is load-bearing on its own -- every
# headline number was produced in one and checked in another.
_ts_idx = None
for _i, _sl in enumerate(S):
    for _sh in _sl.shapes:
        if _sh.has_text_frame and _sh.text_frame.text.strip().startswith("Timeline, Milestones"):
            _ts_idx = _i; break
    if _ts_idx is not None:
        break
if _ts_idx is None:
    raise SystemExit("timeline slide not found for tech-stack insertion")

s_ts = clone_after(p, _ts_idx, _ts_idx + 1)
set_title(s_ts, "Tools \u2014 and what each one produced")
for _sh in list(s_ts.shapes):
    if _sh.has_text_frame and _sh != find_shape(s_ts, "Timeline & milestones"):
        pass
# strip the cloned body, keep title/logo/page number
for _sh in list(s_ts.shapes):
    if _sh.has_text_frame:
        _t = _sh.text_frame.text.strip()
        if _t and not _t.isdigit() and not _t.startswith("Tools \u2014"):
            _sh._element.getparent().remove(_sh._element)

TECH = [
    ("ngspice 42", "every transient in the study",
     "34,622 simulations \u00b7 1.65 V spurious \u00b7 2.58 V margin \u00b7 5.2 % ceiling"),
    ("LTspice 24", "independent re-run of the shipped netlists",
     "1.6487 / 0.8282 / \u22121.1768 V \u2014 matches ngspice within 2 mV"),
    ("MATLAB Online", "independent re-analysis of the same CSVs",
     "504 of 720 feasible \u00b7 7-word Pareto front \u00b7 0.039 \u00b5J per point"),
    ("GNU Octave 11.3", "second run of the same MATLAB scripts",
     "agrees with MATLAB to the last printed digit"),
    ("Icarus Verilog", "RTL verification and mutation testing",
     "8 asserted properties pass \u00b7 injected shoot-through caught 221\u00d7"),
    ("Xilinx Vivado 2024.1.2", "FPGA synthesis and timing on xc7a35t",
     "20 LUTs, 20 flip-flops \u00b7 200 MHz met, 1.996 ns slack"),
    ("Python \u00b7 NumPy", "sweep orchestration and the decomposition",
     "720-word search \u00b7 (A) 25.1 % vs (B) 3.9 %"),
]
_y = 1.55
for _name, _did, _got in TECH:
    add_text(s_ts, 0.70, _y, 3.05, 0.42, [
        para([(_name, True)], level=0, sz=1350, spc=0, bullet=False)])
    add_text(s_ts, 3.85, _y, 3.55, 0.42, [
        para([(_did, False)], level=0, sz=1200, spc=0, bullet=False)])
    add_text(s_ts, 7.50, _y, 5.20, 0.42, [
        para([(_got, True)], level=0, sz=1200, spc=0, bullet=False)])
    _y += 0.62
add_text(s_ts, 0.70, _y + 0.18, 12.00, 0.60, [
    para([("No number rests on one tool. ", True),
          ("Every headline figure was produced in one program and checked in another \u2014 "
           "SPICE against SPICE, MATLAB against Octave, yosys against Vivado.", False)],
         level=0, sz=1250, spc=0, bullet=False)])

for n, sl in enumerate(S, start=1):
    renumber(sl, n)
p.save(OUT)
print("tech-stack slide inserted; deck now %d slides" % len(S))

for n, sl in enumerate(S, start=1):
    renumber(sl, n)
p.save(OUT)
print("MATLAB results slide inserted; deck now %d slides" % len(S))

# ================= THE GAP, on its own slide ==============================
# The panel asked for the research gap stated plainly and separately, not as
# the last bullet of the problem slide. Placed straight after the literature
# so it reads: what exists -> what is missing -> what we propose.
s_gap = clone_after(p, len(S) - 1, 6)
for sh in list(s_gap.shapes):
    if sh.has_text_frame and sh.text_frame.text.strip() and sh.top and sh.top > Inches(1.0):
        sh._element.getparent().remove(sh._element)
    elif sh.shape_type == 13 and sh.left is not None and sh.left < Inches(11):
        sh._element.getparent().remove(sh._element)
S = list(p.slides)
# The clone source carries no usable title placeholder, so place the title and
# the page-number box explicitly, matching the geometry the other slides use.
add_text(s_gap, 0.60, 0.35, 10.50, 0.75,
         [para([("The gap this project fills", True)], level=0, sz=2800, spc=0, bullet=False)])
add_text(s_gap, 12.53, 7.05, 0.50, 0.30,
         [para([("7", False)], level=0, sz=1100, spc=0, bullet=False)])

add_text(s_gap, 0.70, 1.35, 12.10, 1.00, [
    para([("Every active-gate-driver paper reports ONE number: the improvement over a "
           "conventional driver.", True)], level=0, sz=1500, spc=120, bullet=False),
    para([("That number bundles two effects that cost completely different hardware.",
           False)], level=0, sz=1300, spc=0, bullet=False)])

GAPBOX = [
 ("Effect 1  —  a better FIXED setting",
  "Pick a better control word once, at design time. Costs nothing at run time: "
  "no sensor, no ADC, no lookup table, no controller.", "25.1 %", "of baseline"),
 ("Effect 2  —  ADAPTING per operating point",
  "Change the word as load, bus voltage and temperature move. This is what needs "
  "the sensing hardware the architecture is sold on.", "3.9 %", "of baseline"),
]
for i, (head, body, num, unit) in enumerate(GAPBOX):
    x = 0.70 + i * 6.30
    add_text(s_gap, x, 2.60, 5.90, 0.40,
             [para([(head, True)], level=0, sz=1450, spc=0, bullet=False)])
    add_text(s_gap, x, 3.10, 4.05, 1.60,
             [para([(body, False)], level=0, sz=1200, spc=0, bullet=False)])
    stat(s_gap, x + 4.20, 3.05, 1.70, num, unit, "")

add_text(s_gap, 0.70, 4.95, 12.10, 2.10, [
    para([("THE GAP: ", True),
          ("no published work separates them. Nobody reports how much of an active gate "
           "driver's benefit actually requires adaptation, because separating the two needs an "
           "exhaustive search of the control word at every corner — and nobody has run one.",
           False)], level=0, sz=1400, spc=170, bullet=False),
    para([("WHY IT MATTERS: ", True),
          ("only Effect 2 justifies the sense + ADC + lookup table. If it is small, that "
           "hardware is mostly paying for something a design-time choice already delivers \u2014 "
           "a buy-or-not decision nobody has costed.", False)],
         level=0, sz=1400, spc=170, bullet=False),
    para([("WHAT WE DO: ", True),
          ("720 control words × 4 corners, exhaustively, in ngspice — then report the "
           "two numbers separately instead of their sum.", False)],
         level=0, sz=1400, spc=0, bullet=False),
])
for n, sl in enumerate(S, start=1):
    renumber(sl, n)
p.save(OUT)
print("THE GAP slide inserted; deck now %d slides" % len(S))

# ============ BASE PAPER, IMPLEMENTED AND RUN =============================
# The panel asked for the base paper to be BUILT, not just cited, so that
# "what they did / what we add" is a measurement and not an assertion.
s_bp = clone_after(p, len(S) - 1, 12)
for sh in list(s_bp.shapes):
    if sh.has_text_frame and sh.text_frame.text.strip() and sh.top and sh.top > Inches(1.0):
        sh._element.getparent().remove(sh._element)
    elif sh.shape_type == 13 and sh.left is not None and sh.left < Inches(11):
        sh._element.getparent().remove(sh._element)
S = list(p.slides)
add_text(s_bp, 0.60, 0.35, 10.50, 0.75,
         [para([("We implemented the base paper, then ours", True)],
               level=0, sz=2800, spc=0, bullet=False)])
add_text(s_bp, 12.53, 7.05, 0.50, 0.30,
         [para([("13", False)], level=0, sz=1100, spc=0, bullet=False)])
add_text(s_bp, 0.70, 1.32, 12.10, 0.95, [
    para([("Citing a base paper is not a comparison. ", True),
          ("models/basedrv.lib implements Takayama, Okuda & Hikihara's DAC-architecture "
           "driver — a multibit code that changes DURING the switching edge, with no "
           "Miller clamp and no negative off-bias rail, because those are ours. It runs "
           "inside sim/dpt.cir, byte-identical otherwise, so only the driver differs.",
           False)], level=0, sz=1250, spc=0, bullet=False)])
BPC = [
 ("Base paper, as we built it", "+0.533 V", False,
  "Their time-sequenced multibit code alone already clears the 1.4 V threshold."),
 ("Ours, constant code, no clamp", "−0.249 V", False,
  "FALSE TURN-ON. A fast FIXED code is worse than their sequenced one — their "
  "contribution is real, and we reproduce it."),
 ("Ours, active Miller clamp on", "+0.570 V", False,
  "Only marginally past the base paper. The clamp alone is not the story."),
 ("Ours, clamp + −2 V off-bias", "+2.576 V", True,
  "4.8× the base paper's margin. This is the shipped configuration."),
]
for i, (lab, val, hot, note) in enumerate(BPC):
    y = 2.62 + i * 0.70
    add_text(s_bp, 0.70, y, 3.90, 0.58,
             [para([(lab, hot)], level=0, sz=1300, spc=0, bullet=False)])
    add_text(s_bp, 4.75, y - 0.06, 1.55, 0.48,
             [para([(val, True)], level=0, sz=1700, spc=0, bullet=False)])
    add_text(s_bp, 6.50, y, 6.10, 0.60,
             [para([(note, False)], level=0, sz=1150, spc=0, bullet=False)])
add_text(s_bp, 0.70, 5.62, 12.10, 1.35, [
    para([("What the comparison shows. ", True),
          ("The base paper shapes the gate waveform in TIME; we hold the code constant and "
           "add two actuators it does not have. Both clear the threshold — theirs works "
           "— but ours clears it by 4.8× more, and ours is the one whose settings "
           "can then be searched exhaustively to ask what adaptation is actually worth.",
           False)], level=0, sz=1250, spc=150, bullet=False),
    para([("Reproduce: ", True), ("python3 scripts/basepaper_compare.py", True),
          ("  — four ngspice runs, prints this table.", False)],
         level=0, sz=1150, spc=0, bullet=False),
])
# ---------------- SECTION DIVIDERS ----------------------------------------
# Thirty slides of white with navy headings reads as raw material rather than
# a designed deck. Four dark dividers break it into acts and give the eye a
# rest; they cost two seconds each to present.
#
# Palette is deliberately not the default deck-builder blue-purple: VIT navy
# for the ground, a warm off-white for type, a muted slate for the numerals,
# and the same oxblood already used for "failing" in every figure. Nothing
# gradient, nothing decorative -- the contrast does the work.
NAVY   = RGBColor(0x14, 0x23, 0x3D)
PAPER  = RGBColor(0xF4, 0xF2, 0xEE)
SLATE  = RGBColor(0x6B, 0x7C, 0x93)
OXBLD  = RGBColor(0xB0, 0x00, 0x00)

DIVIDERS = [
    ("Problem Statement & Background", "01", "The problem",
     "A real failure, and the question nobody has answered"),
    ("Proposed Solution & Methodology", "02", "What we built",
     "A 720-point control word, and an exhaustive search"),
    ("Result 1", "03", "What we found",
     "Choosing well beats adapting \u2014 and by how much"),
    ("Conclusion & next steps", "04", "Where this goes",
     "The deliverable, the limits, and Review-II"),
]

def _divider_before(title_prefix, numeral, heading, standfirst):
    idx = None
    for i, sl in enumerate(S):
        for sh in sl.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip().startswith(title_prefix):
                idx = i; break
        if idx is not None:
            break
    if idx is None:
        return False
    d = clone_after(p, idx, idx)
    for sh in list(d.shapes):
        sh._element.getparent().remove(sh._element)
    bg = d.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                            Inches(13.333), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background(); bg.shadow.inherit = False
    if bg.has_text_frame:
        bg.text_frame.text = ""

    def dtext(x, y, w, h, text, size, color, bold=False, italic=False):
        tb = d.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        pa = tf.paragraphs[0]; r = pa.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = "Calibri"
        return tb

    # 1.25 not 1.50: at 1.50 this box ran to 3.55 in and overlapped the
    # heading box starting at 3.35 in.
    dtext(1.30, 2.05, 3.0, 1.25, numeral, 66, SLATE, bold=True)
    dtext(1.30, 3.35, 10.6, 1.1, heading, 44, PAPER, bold=True)
    dtext(1.40, 4.55, 10.4, 0.8, standfirst, 17, SLATE)
    return True

for _pfx, _num, _head, _sf in DIVIDERS:
    _divider_before(_pfx, _num, _head, _sf)
for n, sl in enumerate(S, start=1):
    renumber(sl, n)
p.save(OUT)
print("section dividers inserted; deck now %d slides" % len(S))

# ---------------- references: split across two slides ---------------------
# 30 citations need 9.6 in in a 4.75 in box. Splitting is the honest fix: a
# smaller font would fit but nobody can read 7 pt from the back of a room.
ridx = None
for i, sl in enumerate(S):
    for sh in sl.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip().startswith("References"):
            ridx = i; break
    if ridx is not None: break
if ridx is None:
    raise SystemExit("references slide not found")

ref_body = None
for sh in S[ridx].shapes:
    if sh.has_text_frame and "[1]" in sh.text_frame.text:
        ref_body = sh; break

half = (len(R.REFS) + 1) // 2
first, second = R.REFS[:half], R.REFS[half:]

def ref_paras(items, with_note):
    out = []
    for r in items:
        runs = [("[%d]  " % r.n, True)]
        if r.base:
            runs.append(("BASE PAPER — ", True))
        runs.append((r.cite(), False))
        out.append(para(runs, level=0, sz=900, spc=60, bullet=False))
    if with_note:
        out.append(para([("All %d verified against the publisher record — authors, volume, "
                          "issue and pages resolved through Crossref by DOI and title, not "
                          "transcribed by hand." % len(R.REFS), False)],
                        level=0, sz=850, spc=160, bullet=False))
    return out

S2 = clone_after(p, ridx, ridx + 1)
set_title(S[ridx], "References  (1\u2013%d)" % first[-1].n)
set_title(S2, "References  (%d\u2013%d)" % (second[0].n, second[-1].n))
set_body(ref_body, ref_paras(first, False))
for sh in S2.shapes:
    if sh.has_text_frame and "[1]" in sh.text_frame.text:
        set_body(sh, ref_paras(second, True)); break
p.save(OUT)
print("references split: %d + %d across two slides" % (len(first), len(second)))

# ---------------- review date --------------------------------------------
# template_ext.pptx carries 02.09.2026 on slides 2 and 3. The review moved to
# the 9th. Both are Wednesdays, so the day name on slide 2 stays correct.
REVIEW_DATE_OLD, REVIEW_DATE_NEW = "02.09.2026", "09.09.2026"
dhits = 0
for sl in S:
    for sh in sl.shapes:
        if not sh.has_text_frame:
            continue
        for pa in sh.text_frame.paragraphs:
            for r in pa.runs:
                if REVIEW_DATE_OLD in r.text:
                    r.text = r.text.replace(REVIEW_DATE_OLD, REVIEW_DATE_NEW)
                    dhits += 1
print("review date set to %s in %d run(s)" % (REVIEW_DATE_NEW, dhits))

for n, sl in enumerate(list(p.slides), start=1):
    renumber(sl, n)
    for _sh in sl.shapes:
        if _sh.has_text_frame and "Use IEEE format" in _sh.text_frame.text:
            for _pa in _sh.text_frame.paragraphs:
                for _r in _pa.runs:
                    _r.text = ""
p.save(OUT)
print("base-paper implementation slide inserted; deck now %d slides" % len(list(p.slides)))
