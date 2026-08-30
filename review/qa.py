# -*- coding: utf-8 -*-
"""Geometry + text-fit QA. Substitutes for visual QA, which is unavailable
here: LibreOffice in this container fails to load even a trivial HTML file,
so no render can be produced. Estimates use a Calibri advance of 0.50 em
(true value ~0.478), so the checker errs toward reporting overflow."""
from pptx import Presentation
from pptx.util import Emu
EMU = 914400.0
CH_EM, LINE = 0.50, 1.22
SW, SH = 13.333, 7.5
MARGIN = 0.5

p = Presentation("Review1_GaN_Segmented_Gate_Driver.pptx")
problems = []
for idx, s in enumerate(p.slides, 1):
    boxes = []
    for sh in s.shapes:
        if sh.left is None or sh.top is None or not sh.width or not sh.height:
            continue
        x, y = sh.left/EMU, sh.top/EMU
        w, h = sh.width/EMU, sh.height/EMU
        name = (sh.name or "")[:20]
        boxes.append((x, y, w, h, name, sh))
        # slide bounds
        if x < -0.01 or y < -0.01 or x+w > SW+0.01 or y+h > SH+0.01:
            problems.append("s%d %-20s outside slide: x=%.2f y=%.2f w=%.2f h=%.2f"
                            % (idx, name, x, y, w, h))
        # text fit
        if sh.has_text_frame and sh.text_frame.text.strip():
            need = 0.0
            for pa in sh.text_frame.paragraphs:
                txt = "".join(r.text for r in pa.runs)
                if not txt:
                    continue
                szs = [r.font.size.pt for r in pa.runs if r.font.size]
                sz = max(szs) if szs else 16.0
                pPr = pa._pPr
                marL = pPr.get("marL") if pPr is not None else None
                ind = (int(marL)/EMU) if marL else 0.0
                avail = max(0.5, w - ind)
                cw = sz*CH_EM/72.0
                lines = max(1, int(len(txt)*cw/avail) + (1 if (len(txt)*cw) % avail else 0))
                spc = 0.0
                if pPr is not None:
                    sa = pPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}spcAft")
                    if sa is not None and len(sa):
                        v = sa[0].get("val")
                        if v: spc = int(v)/100.0/72.0
                need += lines*sz*LINE/72.0 + spc
            if need > h + 0.02:
                problems.append("s%d %-20s TEXT OVERFLOW: needs %.2f in, box %.2f in"
                                % (idx, name, need, h))
    # pairwise overlap among the shapes we placed (skip logo/page-number)
    for i in range(len(boxes)):
        for j in range(i+1, len(boxes)):
            a, b = boxes[i], boxes[j]
            ox = min(a[0]+a[2], b[0]+b[2]) - max(a[0], b[0])
            oy = min(a[1]+a[3], b[1]+b[3]) - max(a[1], b[1])
            if ox > 0.05 and oy > 0.05:
                problems.append("s%d OVERLAP %.2fx%.2f in: %-18s / %s"
                                % (idx, ox, oy, a[4], b[4]))
print("\n".join(problems) if problems else "no geometry or text-fit problems found")
print("\n--- picture placements ---")
for idx, s in enumerate(p.slides, 1):
    for sh in s.shapes:
        if sh.shape_type == 13 and sh.width and sh.width/EMU > 3:
            print("s%d image x=%.2f y=%.2f w=%.2f h=%.2f  bottom=%.2f"
                  % (idx, sh.left/EMU, sh.top/EMU, sh.width/EMU, sh.height/EMU,
                     (sh.top+sh.height)/EMU))
